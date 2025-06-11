#!/usr/bin/env python3
"""
PowerPoint to Slidev Converter
Converts PowerPoint (.pptx) files exported from Google Slides to Slidev presentations
"""

import re
import json
import os
from pathlib import Path
from typing import List, Dict, Any
from dataclasses import dataclass

try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️  python-pptx not installed. Install with: pip install python-pptx")

@dataclass
class SlideContent:
    title: str
    content: List[str]
    slide_type: str = "default"
    layout: str = "default"
    notes: str = ""
    images: List[str] = None

    def __post_init__(self):
        if self.images is None:
            self.images = []

class SlidevConverter:
    def __init__(self):
        self.slides: List[SlideContent] = []
        
    def extract_from_powerpoint(self, pptx_path: str) -> List[SlideContent]:
        """Extract content from PowerPoint file"""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx library required for PowerPoint processing")
        
        presentation = Presentation(pptx_path)
        slides = []
        
        for slide_num, slide in enumerate(presentation.slides):
            slide_content = self.parse_powerpoint_slide(slide, slide_num)
            if slide_content:
                slides.append(slide_content)
        
        return slides
    
    def parse_powerpoint_slide(self, slide, slide_num: int) -> SlideContent:
        """Parse individual PowerPoint slide"""
        title = ""
        content = []
        notes = ""
        images = []
        
        # Extract title from title placeholder
        if slide.shapes.title:
            # Clean up title by replacing newlines with spaces and normalizing whitespace
            title = ' '.join(slide.shapes.title.text.strip().split())
            # Escape angle brackets to prevent Vue parsing errors
            title = title.replace('<', '&lt;').replace('>', '&gt;')
        
        # Extract content from text shapes and images
        for shape in slide.shapes:
            # Handle images
            if hasattr(shape, 'image'):
                try:
                    # Generate a filename for the image
                    image_filename = f"slide_{slide_num + 1}_image_{len(images) + 1}.png"
                    images.append({
                        'filename': image_filename,
                        'image_data': shape.image.blob
                    })
                except Exception as e:
                    print(f"Warning: Could not extract image from slide {slide_num + 1}: {e}")
            
            elif hasattr(shape, "text_frame") and shape.text_frame:
                # Skip if this is the title
                if shape == slide.shapes.title:
                    continue
                
                # Process each paragraph in the text frame
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        # Escape angle brackets to prevent Vue parsing errors
                        text = text.replace('<', '&lt;').replace('>', '&gt;')
                        
                        # Get the indentation level (0 = main bullet, 1 = sub-bullet, etc.)
                        level = paragraph.level if hasattr(paragraph, 'level') else 0
                        
                        # Store text with level information
                        content.append((level, text))
            elif hasattr(shape, "text") and shape.text.strip():
                # Fallback for shapes without text_frame
                text = shape.text.strip()
                
                # Skip if this is the title
                if shape == slide.shapes.title:
                    continue
                
                # Parse bullet points and paragraphs
                if '\n' in text:
                    lines = [line.strip() for line in text.split('\n') if line.strip()]
                    # Escape angle brackets to prevent Vue parsing errors
                    for line in lines:
                        line = line.replace('<', '&lt;').replace('>', '&gt;')
                        content.append((0, line))  # Default to level 0
                else:
                    # Escape angle brackets to prevent Vue parsing errors
                    text = text.replace('<', '&lt;').replace('>', '&gt;')
                    content.append((0, text))  # Default to level 0
        
        # Extract speaker notes
        if slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        
        # Determine slide type
        slide_type = "default"
        if slide_num == 0:
            slide_type = "title"
        elif not content and title:
            slide_type = "section"
        elif len(content) > 3:
            slide_type = "bullets"
        
        # Store image data temporarily for saving later
        self._temp_images = getattr(self, '_temp_images', [])
        self._temp_images.extend(images)
        
        return SlideContent(
            title=title or f"Slide {slide_num + 1}",
            content=content,
            slide_type=slide_type,
            notes=notes,
            images=[img['filename'] for img in images]
        )
        
    def parse_google_doc_content(self, content: str) -> List[SlideContent]:
        """Parse Google Docs content and extract slide structure"""
        lines = content.split('\n')
        current_slide = None
        slides = []
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # Main title (# header)
            if line.startswith('# '):
                if current_slide:
                    slides.append(current_slide)
                current_slide = SlideContent(
                    title=line[2:],
                    content=[],
                    slide_type="title"
                )
            
            # Section headers (## or ###)
            elif line.startswith('## ') or line.startswith('### '):
                if current_slide:
                    slides.append(current_slide)
                current_slide = SlideContent(
                    title=line.replace('#', '').strip(),
                    content=[],
                    slide_type="section"
                )
            
            # Bullet points
            elif line.startswith('* ') or line.startswith('- '):
                if current_slide:
                    current_slide.content.append(line)
                else:
                    # Create a new slide for orphaned bullets
                    current_slide = SlideContent(
                        title="Content",
                        content=[line],
                        slide_type="bullets"
                    )
        
        if current_slide:
            slides.append(current_slide)
            
        return slides
    
    def generate_frontmatter(self, title: str, author: str = "Kenneth Kousen") -> str:
        """Generate Slidev frontmatter"""
        # Escape title for YAML - wrap in quotes to handle special characters and spaces
        safe_title = f'"{title}"'
        return f"""---
theme: seriph
background: https://source.unsplash.com/1920x1080/?java,programming
class: text-center
highlighter: shiki
lineNumbers: false
info: |
  ## {title}
  
  By {author}
  
  Learn more at [KouseniT](https://kousenit.com)
drawings:
  persist: false
transition: slide-left
title: {safe_title}
mdc: true
---"""

    def convert_slide_to_markdown(self, slide: SlideContent) -> str:
        """Convert a slide to Slidev markdown format"""
        if slide.slide_type == "title":
            return f"""
# {slide.title}

<div class="pt-12">
  <span @click="$slidev.nav.next" class="px-2 py-1 rounded cursor-pointer" hover="bg-white bg-opacity-10">
    Press Space for next page <carbon:arrow-right class="inline"/>
  </span>
</div>
"""
        
        elif slide.slide_type == "section":
            # Convert bullet points to v-clicks for animations
            content_md = ""
            if slide.content:
                content_md = "\n<v-clicks>\n\n"
                for item in slide.content:
                    # Handle both tuple format (level, text) and plain text
                    if isinstance(item, tuple):
                        level, text = item
                        # Add indentation based on level
                        indent = "  " * level
                        content_md += f"{indent}- {text}\n"
                    else:
                        # Fallback for plain text
                        content_md += f"- {item}\n"
                content_md += "\n</v-clicks>"
            
            # Add images if any
            image_md = ""
            if slide.images:
                image_md = "\n\n"
                for image in slide.images:
                    image_md += f'<img src="./{image}" alt="Image" style="max-width: 80%; max-height: 400px; margin: 20px auto; display: block;" />\n\n'
            
            return f"""
# {slide.title}

{content_md}{image_md}
"""
        
        else:  # default bullets
            content_md = "\n<v-clicks>\n\n"
            for item in slide.content:
                # Handle both tuple format (level, text) and plain text
                if isinstance(item, tuple):
                    level, text = item
                    # Add indentation based on level
                    indent = "  " * level
                    content_md += f"{indent}- {text}\n"
                else:
                    # Fallback for plain text
                    content_md += f"- {item}\n"
            content_md += "\n</v-clicks>"
            
            # Add images if any
            image_md = ""
            if slide.images:
                image_md = "\n\n"
                for image in slide.images:
                    image_md += f'<img src="./{image}" alt="Image" style="max-width: 80%; max-height: 400px; margin: 20px auto; display: block;" />\n\n'
            
            return f"""
# {slide.title}

{content_md}{image_md}
"""

    def add_closing_slide(self, author: str = "Kenneth Kousen") -> str:
        """Generate a closing thank you slide"""
        return f"""
# Thank You!

<div class="text-center">

## Questions?

<div class="pt-12">
  <span class="text-6xl"><carbon:logo-github /></span>
</div>

**{author}**  
*Author, Speaker, Java & AI Expert*

[kousenit.com](https://kousenit.com) | [@kenkousen](https://twitter.com/kenkousen)

</div>
"""

    def convert_powerpoint_to_slidev(self, pptx_path: str, output_dir: str = ".", 
                                   author: str = "Kenneth Kousen") -> str:
        """Convert PowerPoint file to Slidev presentation"""
        
        # Initialize temp images storage
        self._temp_images = []
        
        # Extract slides from PowerPoint
        slides = self.extract_from_powerpoint(pptx_path)
        
        # Get title from first slide or filename
        title = slides[0].title if slides else Path(pptx_path).stem
        
        # Create presentation-specific subdirectory
        presentation_dir = self.sanitize_filename(title)
        output_path_obj = Path(output_dir) / presentation_dir
        output_path_obj.mkdir(parents=True, exist_ok=True)
        
        # Save extracted images
        if hasattr(self, '_temp_images') and self._temp_images:
            print(f"💾 Saving {len(self._temp_images)} images...")
            for image_data in self._temp_images:
                try:
                    image_path = output_path_obj / image_data['filename']
                    with open(image_path, 'wb') as f:
                        f.write(image_data['image_data'])
                    print(f"   ✅ Saved {image_data['filename']}")
                except Exception as e:
                    print(f"   ❌ Failed to save {image_data['filename']}: {e}")
        
        # Generate Slidev markdown
        slidev_content = []
        
        # Add frontmatter
        slidev_content.append(self.generate_frontmatter(title, author))
        
        # Convert each slide
        for i, slide in enumerate(slides):
            # Skip the separator before the first slide (title slide)
            # since it should be part of the frontmatter
            if i > 0:
                slidev_content.append("---")
            slidev_content.append(self.convert_slide_to_markdown(slide))
        
        # Add closing slide
        slidev_content.append("---")
        slidev_content.append(self.add_closing_slide(author))
        
        # Join everything
        full_content = "\n".join(slidev_content)
        
        # Write to file as slides.md (Slidev default)
        output_path = output_path_obj / "slides.md"
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(full_content)
        
        print(f"✅ Slidev presentation created: {output_path}")
        print(f"📊 Converted {len(slides)} slides")
        print(f"📁 Output directory: {output_path_obj}")
        return str(output_path)
    
    def sanitize_filename(self, title: str) -> str:
        """Convert title to valid filename"""
        # Remove special characters and replace spaces with hyphens
        clean = re.sub(r'[^\w\s-]', '', title)
        clean = re.sub(r'[-\s]+', '-', clean)
        return clean.lower()

# MCP Integration Class
class MCPSlidevConverter:
    def __init__(self):
        self.converter = SlidevConverter()
    
    def convert_from_powerpoint(self, pptx_path: str, 
                              output_dir: str = "./slidev-presentations") -> str:
        """Convert PowerPoint file to Slidev"""
        return self.converter.convert_powerpoint_to_slidev(pptx_path, output_dir)
    
    def convert_from_google_drive_doc(self, doc_content: str, title: str, 
                                    output_dir: str = "./slidev-presentations") -> str:
        """Convert Google Docs content fetched via MCP to Slidev"""
        return self.converter.convert_to_slidev(doc_content, title, output_dir)
    
    def batch_convert_powerpoints(self, pptx_directory: str, 
                                output_dir: str = "./slidev-presentations") -> List[str]:
        """Convert multiple PowerPoint files in batch"""
        results = []
        pptx_files = Path(pptx_directory).glob("*.pptx")
        
        for pptx_file in pptx_files:
            try:
                result = self.convert_from_powerpoint(str(pptx_file), output_dir)
                results.append(result)
            except Exception as e:
                print(f"❌ Error converting {pptx_file.name}: {e}")
        
        return results

# CLI Usage Example
def main():
    """Example usage of the converter"""
    import argparse
    
    parser = argparse.ArgumentParser(description="Convert presentations to Slidev format")
    parser.add_argument("--pptx", help="PowerPoint file to convert")
    parser.add_argument("--batch", help="Directory containing PowerPoint files")
    parser.add_argument("--output", default="./slidev-presentations", help="Output directory")
    
    args = parser.parse_args()
    
    converter = MCPSlidevConverter()
    
    if args.pptx:
        # Convert single PowerPoint file
        result = converter.convert_from_powerpoint(args.pptx, args.output)
        print(f"✅ Conversion complete! File saved to: {result}")
        
    elif args.batch:
        # Batch convert PowerPoint files
        results = converter.batch_convert_powerpoints(args.batch, args.output)
        print(f"✅ Batch conversion complete! Converted {len(results)} files")
        
    else:
        # Demo with sample content
        sample_content = """
# ESSENTIAL Generative AI Tips ALL Java Developers NEED To Master in 2025

## Spring AI
* If you're invested in the Spring framework, this has everything you need
* Use configuration properties to configure multiple AI models

## LangChain4j  
* Define and use AI services. They rock.
* The langchain4j-examples repo has a sample for everything
"""
        
        result = converter.convert_from_google_drive_doc(
            sample_content, 
            "Essential AI Tips for Java Developers 2025",
            args.output
        )
        print(f"✅ Demo conversion complete! File saved to: {result}")
    
    print("\n🚀 To view your presentation:")
    print("1. cd [output-directory]/[presentation-name]/")
    print("2. npm install -g @slidev/cli")
    print("3. slidev")

if __name__ == "__main__":
    main()
"""

## Integration with Claude

This script can be enhanced to work with Claude by:

1. **MCP Integration**: Use the Google Drive MCP server to fetch documents
2. **Batch Processing**: Process multiple presentations at once
3. **Template Customization**: Allow different Slidev themes and layouts
4. **Asset Management**: Handle images and media files
5. **Preview Generation**: Automatically generate preview images

## Usage with your presentations:

```bash
# Using Claude with this script
claude run slides_converter.py --pptx "presentation.pptx"

# Or batch convert multiple presentations
claude run slides_converter.py --batch ./presentations/
```
"""
